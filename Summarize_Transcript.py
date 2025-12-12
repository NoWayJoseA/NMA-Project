import io
import os
import re
import csv
import time
import argparse
import subprocess
from typing import Tuple, Optional, Iterable

# Google Drive API
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from googleapiclient.errors import HttpError
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request

# Content extraction
from pdfminer.high_level import extract_text as pdf_extract_text
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation

# OpenAI (Responses + Audio APIs)
from openai import OpenAI

# -------- Settings --------
# If you need broader Drive access (like write, delete, or sharing), this scope will need to change.
# For this project we only read files, so read-only is safer.
SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]

# This CSV is what the website reads (search page).
# If you change the filename or move it into a subfolder, make sure the front-end CSV_URL gets updated too.
OUTPUT_CSV = "drive_summaries.csv"

# Small, fast model for summarizing text/transcripts.
# If you swap to a bigger model (e.g. gpt-4.1), update the default here.
OPENAI_MODEL = os.getenv("OPENAI_MODEL", "gpt-4o-mini")

# Summarization prompt
SYSTEM_PROMPT = (
    "You are a concise assistant. Summarize the document in 3–4 sentences. "
    "Mention the main topic, key points, and any calls-to-action. Avoid fluff."
)

# Google MIME types
GOOGLE_DOC = "application/vnd.google-apps.document"
GOOGLE_SHEET = "application/vnd.google-apps.spreadsheet"
GOOGLE_SLIDE = "application/vnd.google-apps.presentation"
GOOGLE_FOLDER = "application/vnd.google-apps.folder"
GOOGLE_SHORTCUT = "application/vnd.google-apps.shortcut"


# -------------------------- Auth & Service --------------------------

def get_drive_service():
    """
    Spins up the Drive service. Reuses token.json if it exists so you don't have to re-auth every run.
    - token.json holds your refresh token, so don't commit it to GitHub.
    - credentials.json is the OAuth client from Google Cloud (Desktop app).
    """
    creds = None
    if os.path.exists("token.json"):
        creds = Credentials.from_authorized_user_file("token.json", SCOPES)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            # Token is tired, give it a coffee (refresh)
            creds.refresh(Request())
        else:
            # First-time auth dance
            if not os.path.exists("credentials.json"):
                raise FileNotFoundError(
                    "credentials.json not found. Create an OAuth 2.0 Desktop Client in "
                    "Google Cloud → download JSON → save as credentials.json next to this script."
                )
            flow = InstalledAppFlow.from_client_secrets_file("credentials.json", SCOPES)
            # WSL-friendly: prints URL; open it manually in your browser if needed
            creds = flow.run_local_server(port=0, open_browser=False)

        with open("token.json", "w") as token:
            token.write(creds.to_json())

    return build("drive", "v3", credentials=creds)


# -------------------------- Helpers --------------------------

def parse_folder_id(folder_arg: str) -> str:
    """
    Accepts a folder ID or a full URL and returns the folder ID.
    Handles the usual /folders/<ID> style links.
    This is the CLI input you'll pass via --folder.
    """
    if re.match(r"^[A-Za-z0-9_\-]{10,}$", folder_arg):
        return folder_arg
    m = re.search(r"/folders/([A-Za-z0-9_\-]+)", folder_arg)
    if m:
        return m.group(1)
    raise ValueError("Could not parse folder ID from input. Provide a valid folder ID or folder URL.")


def list_folder_contents(service, folder_id: str) -> Iterable[dict]:
    """
    Lists direct children of a folder (files + subfolders). We recurse elsewhere.
    If you want to filter out certain types at the query level, this is the place.
    """
    page_token = None
    # Asking Drive for the stuff we actually care about
    fields = (
        "nextPageToken, files("
        "id, name, mimeType, webViewLink, modifiedTime, createdTime, "
        "parents, owners(displayName,emailAddress), shortcutDetails)"
    )
    query = f"'{folder_id}' in parents and trashed = false"

    while True:
        resp = service.files().list(
            q=query,
            fields=fields,
            pageSize=1000,
            pageToken=page_token,
            supportsAllDrives=True,
            includeItemsFromAllDrives=True,
            corpora="allDrives",
        ).execute()
        for f in resp.get("files", []):
            yield f
        page_token = resp.get("nextPageToken")
        if not page_token:
            break


def resolve_shortcut(f: dict) -> dict:
    """
    If item is a Drive shortcut, pretend it's actually the target file.
    This keeps the rest of the code from having to care about shortcuts.
    """
    if f.get("mimeType") != GOOGLE_SHORTCUT:
        return f
    sd = f.get("shortcutDetails", {}) or {}
    target_id = sd.get("targetId")
    target_mime = sd.get("targetMimeType")
    if not target_id or not target_mime:
        return f
    g = dict(f)
    g["id"] = target_id
    g["mimeType"] = target_mime
    return g


def walk_folder(service, root_folder_id: str, recursive: bool = True) -> Iterable[dict]:
    """
    Simple DFS over folders. Yields all files under root.
    If you want to limit depth or skip certain paths, you can hook that here.
    """
    stack = [root_folder_id]
    seen = set()
    while stack:
        current = stack.pop()
        for item in list_folder_contents(service, current):
            item = resolve_shortcut(item)
            if item["mimeType"] == GOOGLE_FOLDER:
                if recursive and item["id"] not in seen:
                    seen.add(item["id"])
                    stack.append(item["id"])
                # we care about files, not folders, so skip summarizing folders
                continue
            yield item


# -------- File download/export helpers (Drive v3) --------

def export_google_doc_to_text(service, file_id) -> str:
    """
    Google Docs → plain text.
    If you want HTML instead, change mimeType to text/html.
    """
    buf = io.BytesIO()
    request = service.files().export(fileId=file_id, mimeType="text/plain")
    buf.write(request.execute())
    return buf.getvalue().decode("utf-8", errors="ignore")


def export_google_sheet_to_text(service, file_id) -> str:
    """
    Google Sheets → CSV.
    This is mostly for giving the model something to chew on; not meant to be pretty.
    """
    request = service.files().export(fileId=file_id, mimeType="text/csv")
    data = request.execute()
    return data.decode("utf-8", errors="ignore")


def export_google_slide_to_pdf_then_text(service, file_id) -> str:
    """
    Google Slides → PDF → run through pdfminer.
    Not perfect, but better than nothing for slide decks.
    """
    request = service.files().export(fileId=file_id, mimeType="application/pdf")
    pdf_bytes = request.execute()
    tmp_path = f"/tmp/{file_id}.pdf"
    with open(tmp_path, "wb") as f:
        f.write(pdf_bytes)
    try:
        return pdf_extract_text(tmp_path)
    finally:
        try:
            os.remove(tmp_path)
        except OSError:
            pass


def download_binary(service, file_id) -> bytes:
    """
    Generic download helper. Works for PDFs, PPTX, audio, video, etc.
    """
    request = service.files().get_media(fileId=file_id)
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return fh.getvalue()


# -------- Local text extraction helpers --------

def text_from_pdf_bytes(b: bytes) -> str:
    p = f"/tmp/_dl_{int(time.time() * 1000)}.pdf"
    with open(p, "wb") as f:
        f.write(b)
    try:
        return pdf_extract_text(p)
    finally:
        try:
            os.remove(p)
        except OSError:
            pass


def text_from_docx_bytes(b: bytes) -> str:
    p = f"/tmp/_dl_{int(time.time() * 1000)}.docx"
    with open(p, "wb") as f:
        f.write(b)
    try:
        doc = DocxDocument(p)
        return "\n".join([para.text for para in doc.paragraphs])
    finally:
        try:
            os.remove(p)
        except OSError:
            pass


def text_from_pptx_bytes(b: bytes) -> str:
    p = f"/tmp/_dl_{int(time.time() * 1000)}.pptx"
    with open(p, "wb") as f:
        f.write(b)
    try:
        prs = PptxPresentation(p)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    chunks.append(shape.text)
        return "\n".join(chunks)
    finally:
        try:
            os.remove(p)
        except OSError:
            pass


# -------------------------- Media helpers (audio/video) --------------------------

def is_media_file(f: dict) -> bool:
    """
    Quick-and-dirty check: is this file audio or video?
    We look at both mimeType and file extension because Drive can be weird.
    """
    mime = f.get("mimeType", "") or ""
    name = (f.get("name") or "").lower()

    if mime.startswith("audio/") or mime.startswith("video/"):
        return True

    media_exts = (
        ".mp3", ".wav", ".m4a", ".aac", ".flac", ".ogg",
        ".mp4", ".mov", ".mkv", ".webm"
    )
    return any(name.endswith(ext) for ext in media_exts)


def write_temp_file(b: bytes, suffix: str) -> str:
    """
    Drops bytes into /tmp with a random-ish name and returns the path.
    """
    path = f"/tmp/_media_{int(time.time() * 1000)}{suffix}"
    with open(path, "wb") as f:
        f.write(b)
    return path


def convert_to_mp3(input_path: str) -> str:
    """
    Uses ffmpeg to turn whatever we have into MP3.
    Requires ffmpeg installed and on PATH.
    """
    output_path = input_path + ".mp3"
    # ffmpeg -y means "yes overwrite" so it doesn't nag
    cmd = ["ffmpeg", "-y", "-i", input_path, "-vn", "-acodec", "libmp3lame", output_path]
    subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    return output_path


def transcribe_media_bytes(client: OpenAI, b: bytes, filename: str) -> str:
    """
    Full pipeline for audio/video:
    - write bytes to /tmp
    - convert to mp3 if it's not already in a friendly format
    - send to gpt-4o-transcribe
    - return raw transcript text

    If you swap transcription models later, change the `model=` value below.
    """
    name_lower = (filename or "").lower()
    # Guess a suffix so ffmpeg doesn't freak out
    if "." in name_lower:
        suffix = "." + name_lower.split(".")[-1]
    else:
        suffix = ".bin"

    input_path = write_temp_file(b, suffix)

    try:
        # If it's already an audio format gpt-4o-transcribe likes, we can skip conversion.
        # Otherwise we just run it through ffmpeg and force mp3.
        audio_path = input_path
        friendly_exts = (".mp3", ".wav", ".m4a", ".webm")
        if not name_lower.endswith(friendly_exts):
            audio_path = convert_to_mp3(input_path)

        with open(audio_path, "rb") as audio_file:
            resp = client.audio.transcriptions.create(
                model="gpt-4o-transcribe",
                file=audio_file,
            )

        # gpt-4o-transcribe returns .text directly, super convenient
        return resp.text.strip()
    finally:
        # Try to clean up after ourselves. If this fails, it's not the end of the world.
        for path in {input_path, input_path + ".mp3"}:
            try:
                if os.path.exists(path):
                    os.remove(path)
            except OSError:
                pass


# -------------------------- Summarization (OpenAI Responses API) --------------------------

def summarize_text(client: OpenAI, title: str, text: str) -> str:
    """
    Classic "stuff a bunch of text into a small summary" helper.
    If you need more/less detail, tweak SYSTEM_PROMPT or MAX_CHARS.
    """
    MAX_CHARS = 120_000  # big enough for long docs but not insane
    body = (text or "")[:MAX_CHARS]
    if not body.strip():
        return ""

    prompt = (
        f"{SYSTEM_PROMPT}\n\n"
        f"Title: {title}\n\n"
        f"Content:\n{body}\n\n"
        f"Return ONLY the 3–4 sentence summary."
    )
    resp = client.responses.create(
        model=OPENAI_MODEL,
        input=prompt,
        temperature=0.2,
    )
    try:
        return resp.output_text.strip()
    except Exception:
        # Fallback if the SDK structure changes on us
        if hasattr(resp, "choices") and resp.choices:
            return getattr(resp.choices[0], "message", {}).get("content", "").strip()
        return ""


def extract_text_for_file(service, f) -> Tuple[str, Optional[str]]:
    """
    Handles "normal" docs: Google Docs, Sheets, Slides, PDFs, Word, PPT, etc.
    Returns (text, strategy_used).
    If you want to special-case certain file types, this is the central place.
    """
    mime = f.get("mimeType", "")
    fid = f["id"]
    try:
        # Native Google file types first
        if mime == GOOGLE_DOC:
            return export_google_doc_to_text(service, fid), "export_doc_txt"
        if mime == GOOGLE_SHEET:
            return export_google_sheet_to_text(service, fid), "export_sheet_csv"
        if mime == GOOGLE_SLIDE:
            return export_google_slide_to_pdf_then_text(service, fid), "export_slide_pdf"

        # Common non-Google formats
        name_lower = f.get("name", "").lower()
        b = download_binary(service, fid)

        if mime == "application/pdf" or name_lower.endswith(".pdf"):
            return text_from_pdf_bytes(b), "pdfminer"

        if mime in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ) or name_lower.endswith((".docx", ".doc")):
            return text_from_docx_bytes(b), "python-docx"

        if mime in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ) or name_lower.endswith((".pptx", ".ppt")):
            return text_from_pptx_bytes(b), "python-pptx"

        # Fallback: just assume it's plain text-ish
        try:
            return b.decode("utf-8", errors="ignore"), "raw_text_decode"
        except Exception:
            return "", "unsupported_type"
    except HttpError as e:
        return f"DOWNLOAD/EXPORT ERROR: {e}", "error"


# -------------------------- Main --------------------------

def main():
    """
    Entry point.
    Things future editors commonly touch:
      - OPENAI_API_KEY (from env)
      - OUTPUT_CSV filename/path
      - which Drive folder to walk (CLI flag: --folder)
      - recursive vs non-recursive crawl
    """

    # OpenAI key check: we don't bother starting Drive auth if we can't talk to OpenAI anyway.
    # IMPORTANT: Set this in your environment, do NOT paste keys into source code.
    if not os.getenv("OPENAI_API_KEY"):
        raise EnvironmentError("Set OPENAI_API_KEY environment variable before running this script.")

    parser = argparse.ArgumentParser(description="Summarize Google Drive files within a specific folder.")
    parser.add_argument("--folder", required=True, help="Google Drive folder ID or URL")
    parser.add_argument("--recursive", dest="recursive", action="store_true", help="Recurse into subfolders (default)")
    parser.add_argument("--no-recursive", dest="recursive", action="store_false", help="Only files directly in the folder")
    parser.set_defaults(recursive=True)
    args = parser.parse_args()

    folder_id = parse_folder_id(args.folder)

    # Client will pick up OPENAI_API_KEY from the environment.
    client = OpenAI()
    drive = get_drive_service()

    # Final CSV columns that the front-end expects.
    # If you add/remove columns here, update the search page parsing and tags there.
    out_fields = [
        "file_id",
        "file_name",
        "file_type",
        "date_uploaded",
        "author",
        "summary",
        "transcript",   # only filled for audio/video files
        "link",
    ]

    seen_ids = set()

    with open(OUTPUT_CSV, "w", newline="", encoding="utf-8") as fh:
        writer = csv.DictWriter(fh, fieldnames=out_fields)
        writer.writeheader()

        for f in walk_folder(drive, folder_id, recursive=args.recursive):
            fid = f["id"]
            if fid in seen_ids:
                continue
            seen_ids.add(fid)

            name = f.get("name", "")
            mime = f.get("mimeType", "")
            owner_info = f.get("owners", [{}])
            owner = (
                owner_info[0].get("displayName")
                or owner_info[0].get("emailAddress")
                or ""
            )
            link = f.get("webViewLink", "")

            # Prefer createdTime for "date uploaded", fall back to modifiedTime if needed
            date_uploaded = f.get("createdTime") or f.get("modifiedTime") or ""

            is_media = is_media_file(f)

            summary = ""
            transcript = ""

            try:
                if is_media:
                    # This is audio or video – we transcribe instead of running pdfminer/docx/etc.
                    print(f"[MEDIA] Downloading + transcribing: {name} [{mime}]")
                    b = download_binary(drive, fid)

                    # Light nap between API calls so we don't get throttled
                    time.sleep(float(os.getenv("PAUSE_BETWEEN_CALLS", "0.25")))

                    transcript = transcribe_media_bytes(client, b, name)

                    # Once we have a transcript, we treat it like any other "text" to summarize.
                    summary = summarize_text(client, name, transcript)
                else:
                    # Regular document path: extract text, then summarize it.
                    print(f"[DOC] Downloading + summarizing: {name} [{mime}]")
                    text, _strategy = extract_text_for_file(drive, f)

                    if text and not text.startswith("DOWNLOAD/EXPORT ERROR"):
                        time.sleep(float(os.getenv("PAUSE_BETWEEN_CALLS", "0.25")))
                        summary = summarize_text(client, name, text)
                    else:
                        summary = text  # if it's an error string, just pass it through
            except Exception as e:
                # If anything blows up, we still log something in the CSV so you can see what went wrong.
                err_msg = f"PROCESSING ERROR: {e}"
                if is_media:
                    transcript = err_msg
                else:
                    summary = err_msg

            row = {
                "file_id": fid,
                "file_name": name,
                "file_type": mime,
                "date_uploaded": date_uploaded,
                "author": owner,
                "summary": summary,
                "transcript": transcript if is_media else "",
                "link": link,
            }
            writer.writerow(row)

            tag = "MEDIA" if is_media else "DOC"
            print(f"Done [{tag}]: {name}")

    print(f"\nDone. Results saved to {OUTPUT_CSV}")


if __name__ == "__main__":
    main()
