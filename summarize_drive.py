import io
import os
import re
import csv
import time
import argparse
from typing import Tuple, Optional, Iterable

# Google Drive API
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from googleapiclient.errors import HttpError
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request

# Content extraction
from pdfminer.high_level import extract_text as pdf_extract_text
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation

# OpenAI (Responses API)
from openai import OpenAI

# -------- Settings --------
SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]
OUTPUT_CSV = "drive_summaries.csv"

# Choose a small, fast model for brief summaries.
OPENAI_MODEL = os.getenv("OPENAI_MODEL", "gpt-4o-mini")  # override if you like

# Summarization prompt
SYSTEM_PROMPT = (
    "You are a concise assistant. Summarize the document in 3–4 sentences. "
    "Mention the main topic, key points, and any calls-to-action. Avoid fluff."
)

# Google MIME types
GOOGLE_DOC = "application/vnd.google-apps.document"
GOOGLE_SHEET = "application/vnd.google-apps.spreadsheet"
GOOGLE_SLIDE = "application/vnd.google-apps.presentation"
GOOGLE_FOLDER = "application/vnd.google-apps.folder"
GOOGLE_SHORTCUT = "application/vnd.google-apps.shortcut"


# -------------------------- Auth & Service --------------------------

def get_drive_service():
    creds = None
    if os.path.exists("token.json"):
        creds = Credentials.from_authorized_user_file("token.json", SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            if not os.path.exists("credentials.json"):
                raise FileNotFoundError(
                    "credentials.json not found. Create an OAuth 2.0 Desktop Client in "
                    "Google Cloud → download JSON → save as credentials.json next to this script."
                )
            flow = InstalledAppFlow.from_client_secrets_file("credentials.json", SCOPES)
            # WSL-friendly: prints URL; open it in your Windows browser
            creds = flow.run_local_server(port=0, open_browser=False)
        with open("token.json", "w") as token:
            token.write(creds.to_json())
    return build("drive", "v3", credentials=creds)


# -------------------------- Helpers --------------------------

def parse_folder_id(folder_arg: str) -> str:
    """
    Accepts a folder ID or a full URL and returns the folder ID.
    Examples:
      https://drive.google.com/drive/folders/<ID>
      https://drive.google.com/drive/u/0/folders/<ID>
    """
    if re.match(r"^[A-Za-z0-9_\-]{10,}$", folder_arg):
        return folder_arg
    m = re.search(r"/folders/([A-Za-z0-9_\-]+)", folder_arg)
    if m:
        return m.group(1)
    raise ValueError("Could not parse folder ID from input. Provide a valid folder ID or folder URL.")


def list_folder_contents(service, folder_id: str) -> Iterable[dict]:
    """
    Lists direct children of a folder (files + subfolders). Does NOT recurse.
    """
    page_token = None
    fields = "nextPageToken, files(id, name, mimeType, webViewLink, modifiedTime, parents, owners(displayName,emailAddress), shortcutDetails)"
    query = f"'{folder_id}' in parents and trashed = false"

    while True:
        resp = service.files().list(
            q=query,
            fields=fields,
            pageSize=1000,
            pageToken=page_token,
            supportsAllDrives=True,
            includeItemsFromAllDrives=True,
            corpora="allDrives",
        ).execute()
        for f in resp.get("files", []):
            yield f
        page_token = resp.get("nextPageToken")
        if not page_token:
            break


def resolve_shortcut(f: dict) -> dict:
    """
    If item is a Drive shortcut, return a pseudo-file dict for its target.
    Otherwise return as-is.
    """
    if f.get("mimeType") != GOOGLE_SHORTCUT:
        return f
    sd = f.get("shortcutDetails", {}) or {}
    target_id = sd.get("targetId")
    target_mime = sd.get("targetMimeType")
    if not target_id or not target_mime:
        return f
    # Keep original metadata but point to target id/mime for processing
    g = dict(f)
    g["id"] = target_id
    g["mimeType"] = target_mime
    return g


def walk_folder(service, root_folder_id: str, recursive: bool = True) -> Iterable[dict]:
    """
    Yields all files under the given folder. If recursive=True, traverses subfolders.
    """
    stack = [root_folder_id]
    seen = set()
    while stack:
        current = stack.pop()
        for item in list_folder_contents(service, current):
            item = resolve_shortcut(item)
            if item["mimeType"] == GOOGLE_FOLDER:
                if recursive and item["id"] not in seen:
                    seen.add(item["id"])
                    stack.append(item["id"])
                # We only summarize files, not the folder itself
                continue
            yield item


# -------- File download/export helpers (Drive v3) --------

def export_google_doc_to_text(service, file_id) -> str:
    buf = io.BytesIO()
    request = service.files().export(fileId=file_id, mimeType="text/plain")
    buf.write(request.execute())
    return buf.getvalue().decode("utf-8", errors="ignore")

def export_google_sheet_to_text(service, file_id) -> str:
    request = service.files().export(fileId=file_id, mimeType="text/csv")
    data = request.execute()
    return data.decode("utf-8", errors="ignore")

def export_google_slide_to_pdf_then_text(service, file_id) -> str:
    request = service.files().export(fileId=file_id, mimeType="application/pdf")
    pdf_bytes = request.execute()
    tmp_path = f"/tmp/{file_id}.pdf"
    with open(tmp_path, "wb") as f:
        f.write(pdf_bytes)
    try:
        return pdf_extract_text(tmp_path)
    finally:
        try:
            os.remove(tmp_path)
        except OSError:
            pass

def download_binary(service, file_id) -> bytes:
    request = service.files().get_media(fileId=file_id)
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return fh.getvalue()

# -------- Local text extraction helpers --------

def text_from_pdf_bytes(b: bytes) -> str:
    p = f"/tmp/_dl_{int(time.time()*1000)}.pdf"
    with open(p, "wb") as f:
        f.write(b)
    try:
        return pdf_extract_text(p)
    finally:
        try:
            os.remove(p)
        except OSError:
            pass

def text_from_docx_bytes(b: bytes) -> str:
    p = f"/tmp/_dl_{int(time.time()*1000)}.docx"
    with open(p, "wb") as f:
        f.write(b)
    try:
        doc = DocxDocument(p)
        return "\n".join([para.text for para in doc.paragraphs])
    finally:
        try:
            os.remove(p)
        except OSError:
            pass

def text_from_pptx_bytes(b: bytes) -> str:
    p = f"/tmp/_dl_{int(time.time()*1000)}.pptx"
    with open(p, "wb") as f:
        f.write(b)
    try:
        prs = PptxPresentation(p)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    chunks.append(shape.text)
        return "\n".join(chunks)
    finally:
        try:
            os.remove(p)
        except OSError:
            pass

# -------- Summarization (OpenAI Responses API) --------

def summarize_text(client: OpenAI, title: str, text: str) -> str:
    MAX_CHARS = 120_000
    body = (text or "")[:MAX_CHARS]
    prompt = (
        f"{SYSTEM_PROMPT}\n\n"
        f"Title: {title}\n\n"
        f"Content:\n{body}\n\n"
        f"Return ONLY the 3–4 sentence summary."
    )
    resp = client.responses.create(
        model=OPENAI_MODEL,
        input=prompt,
        temperature=0.2,
    )
    try:
        return resp.output_text.strip()
    except Exception:
        if hasattr(resp, "choices") and resp.choices:
            return getattr(resp.choices[0], "message", {}).get("content", "").strip()
        return ""

def extract_text_for_file(service, f) -> Tuple[str, Optional[str]]:
    mime = f.get("mimeType", "")
    fid = f["id"]
    try:
        if mime == GOOGLE_DOC:
            return export_google_doc_to_text(service, fid), "export_doc_txt"
        if mime == GOOGLE_SHEET:
            return export_google_sheet_to_text(service, fid), "export_sheet_csv"
        if mime == GOOGLE_SLIDE:
            return export_google_slide_to_pdf_then_text(service, fid), "export_slide_pdf"

        # Common non-Google formats
        name_lower = f.get("name", "").lower()
        b = download_binary(service, fid)

        if mime == "application/pdf" or name_lower.endswith(".pdf"):
            return text_from_pdf_bytes(b), "pdfminer"

        if mime in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ) or name_lower.endswith((".docx", ".doc")):
            return text_from_docx_bytes(b), "python-docx"

        if mime in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ) or name_lower.endswith((".pptx", ".ppt")):
            return text_from_pptx_bytes(b), "python-pptx"

        # Fallback: attempt to decode as text
        try:
            return b.decode("utf-8", errors="ignore"), "raw_text_decode"
        except Exception:
            return "", "unsupported_type"
    except HttpError as e:
        return f"DOWNLOAD/EXPORT ERROR: {e}", "error"


# -------------------------- Main --------------------------

def main():
    # Ensure OpenAI API key exists (do NOT hardcode it in code)
    if not os.getenv("OPENAI_API_KEY"):
        raise EnvironmentError("Set OPENAI_API_KEY environment variable before running.")

    parser = argparse.ArgumentParser(description="Summarize Google Drive files within a specific folder.")
    parser.add_argument("--folder", required=True, help="Google Drive folder ID or URL")
    parser.add_argument("--recursive", dest="recursive", action="store_true", help="Recurse into subfolders (default)")
    parser.add_argument("--no-recursive", dest="recursive", action="store_false", help="Only files directly in the folder")
    parser.set_defaults(recursive=True)
    args = parser.parse_args()

    folder_id = parse_folder_id(args.folder)

    client = OpenAI()
    drive = get_drive_service()

    out_fields = [
        "file_id",
        "name",
        "mime_type",
        "modified_time",
        "owner",
        "webViewLink",
        "extraction_strategy",
        "summary",
    ]

    seen_ids = set()

    with open(OUTPUT_CSV, "w", newline="", encoding="utf-8") as fh:
        writer = csv.DictWriter(fh, fieldnames=out_fields)
        writer.writeheader()

        for f in walk_folder(drive, folder_id, recursive=args.recursive):
            fid = f["id"]
            if fid in seen_ids:
                continue
            seen_ids.add(fid)

            name = f.get("name", "")
            mime = f.get("mimeType", "")

            text, strategy = extract_text_for_file(drive, f)

            summary = ""
            if text and not text.startswith("DOWNLOAD/EXPORT ERROR"):
                try:
                    # gentle spacing to reduce rate issues
                    time.sleep(float(os.getenv("PAUSE_BETWEEN_CALLS", "0.25")))
                    summary = summarize_text(client, name, text)
                except Exception as e:
                    summary = f"SUMMARY ERROR: {e}"

            row = {
                "file_id": fid,
                "name": name,
                "mime_type": mime,
                "modified_time": f.get("modifiedTime"),
                "owner": (f.get("owners", [{}])[0].get("displayName")
                          or f.get("owners", [{}])[0].get("emailAddress") or ""),
                "webViewLink": f.get("webViewLink", ""),
                "extraction_strategy": strategy,
                "summary": summary,
            }
            writer.writerow(row)
            print(f"Processed: {name} [{mime}] -> {strategy}")

    print(f"\nDone. Results saved to {OUTPUT_CSV}")


if __name__ == "__main__":
    main()
